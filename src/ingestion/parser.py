"""Document parsing for multiple file types."""

import csv
import json
import os
import xml.etree.ElementTree as ET
from io import StringIO
from pathlib import Path

from src.utils.logging import get_logger

logger = get_logger(__name__)

# Optional dependencies - don't crash if missing
try:
    from docx import Document as DocxDocument
    _HAS_DOCX = True
except ImportError:
    _HAS_DOCX = False

try:
    import openpyxl
    _HAS_OPENPYXL = True
except ImportError:
    _HAS_OPENPYXL = False

try:
    from pptx import Presentation
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

try:
    from striprtf.striprtf import rtf_to_text
    _HAS_STRIPRTF = True
except ImportError:
    _HAS_STRIPRTF = False

try:
    import pytesseract
    from pdf2image import convert_from_path
    _HAS_OCR = True
except ImportError:
    _HAS_OCR = False

# Supported file extensions and their categories
_TEXT_EXTENSIONS = {".txt", ".text", ".md"}
_HTML_EXTENSIONS = {".html", ".htm"}
_PDF_EXTENSIONS = {".pdf"}
_DOCX_EXTENSIONS = {".docx"}
_XLSX_EXTENSIONS = {".xlsx"}
_CSV_EXTENSIONS = {".csv"}
_PPTX_EXTENSIONS = {".pptx"}
_RTF_EXTENSIONS = {".rtf"}
_JSON_EXTENSIONS = {".json"}
_XML_EXTENSIONS = {".xml"}

ALL_SUPPORTED_EXTENSIONS = (
    _TEXT_EXTENSIONS
    | _HTML_EXTENSIONS
    | _PDF_EXTENSIONS
    | _DOCX_EXTENSIONS
    | _XLSX_EXTENSIONS
    | _CSV_EXTENSIONS
    | _PPTX_EXTENSIONS
    | _RTF_EXTENSIONS
    | _JSON_EXTENSIONS
    | _XML_EXTENSIONS
)


def parse_document(file_path: str) -> list[dict]:
    """Parse a document file and extract text with metadata.

    Returns list of dicts with "text" and "metadata" keys.
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = path.suffix.lower()
    logger.info("Parsing %s (type: %s)", path.name, ext)

    dispatch = {
        **{e: _parse_pdf for e in _PDF_EXTENSIONS},
        **{e: _parse_html for e in _HTML_EXTENSIONS},
        **{e: _parse_text for e in _TEXT_EXTENSIONS},
        **{e: _parse_docx for e in _DOCX_EXTENSIONS},
        **{e: _parse_xlsx for e in _XLSX_EXTENSIONS},
        **{e: _parse_csv for e in _CSV_EXTENSIONS},
        **{e: _parse_pptx for e in _PPTX_EXTENSIONS},
        **{e: _parse_rtf for e in _RTF_EXTENSIONS},
        **{e: _parse_json_file for e in _JSON_EXTENSIONS},
        **{e: _parse_xml for e in _XML_EXTENSIONS},
    }

    parser = dispatch.get(ext)
    if parser:
        return parser(path)

    logger.warning("Unknown file type %s, treating as plain text", ext)
    return _parse_text(path)


def _parse_pdf(path: Path) -> list[dict]:
    """Parse PDF page by page to keep memory bounded on large documents.

    Falls back to OCR if text extraction returns nothing.
    """
    try:
        import fitz  # PyMuPDF - fast, low memory, page-by-page
        return _parse_pdf_pymupdf(path)
    except ImportError:
        pass

    # Fallback to LlamaIndex (loads entire PDF at once, higher memory)
    try:
        from llama_index.core import SimpleDirectoryReader
        reader = SimpleDirectoryReader(input_files=[str(path)])
        documents = reader.load_data()

        results = []
        for i, doc in enumerate(documents):
            text = doc.text.strip()
            if not text:
                continue
            results.append({
                "text": text,
                "metadata": {
                    "source": path.name,
                    "page": str(i + 1),
                    "file_path": str(path),
                },
            })

        if results:
            return results
    except Exception:
        logger.exception("LlamaIndex PDF parsing failed for %s", path.name)

    # If text extraction returned nothing, try OCR
    if _HAS_OCR:
        logger.info("PDF text extraction empty, trying OCR for %s", path.name)
        return _ocr_pdf(path)

    logger.warning("No text extracted from %s and OCR not available", path.name)
    return []


def _parse_pdf_pymupdf(path: Path) -> list[dict]:
    """Parse PDF page by page using PyMuPDF. Much lower memory than LlamaIndex."""
    import fitz

    results = []
    doc = fitz.open(str(path))
    try:
        total_pages = len(doc)
        logger.info("Parsing %d pages from %s", total_pages, path.name)

        for page_num in range(total_pages):
            page = doc[page_num]
            text = page.get_text().strip()
            if not text:
                continue
            results.append({
                "text": text,
                "metadata": {
                    "source": path.name,
                    "page": str(page_num + 1),
                    "total_pages": str(total_pages),
                    "file_path": str(path),
                },
            })
            # Log progress every 50 pages
            if (page_num + 1) % 50 == 0:
                logger.info("  Parsed %d / %d pages", page_num + 1, total_pages)
    finally:
        doc.close()

    # If PyMuPDF got nothing, try OCR
    if not results and _HAS_OCR:
        logger.info("PyMuPDF got no text, trying OCR for %s", path.name)
        return _ocr_pdf(path)

    return results


def _ocr_pdf(path: Path) -> list[dict]:
    """Extract text from a scanned PDF using pytesseract + pdf2image."""
    from config.settings import settings

    if settings.tesseract_cmd:
        pytesseract.pytesseract.tesseract_cmd = settings.tesseract_cmd

    try:
        images = convert_from_path(str(path))
    except Exception:
        logger.warning("pdf2image failed for %s, skipping OCR", path.name)
        return []

    results = []
    for i, img in enumerate(images):
        try:
            text = pytesseract.image_to_string(img)
        except Exception:
            logger.warning("OCR failed on page %d of %s", i + 1, path.name)
            continue
        if text.strip():
            results.append({
                "text": text,
                "metadata": {
                    "source": path.name,
                    "page": str(i + 1),
                    "file_path": str(path),
                    "format": "pdf_ocr",
                },
            })
    return results


def _parse_html(path: Path) -> list[dict]:
    """Parse HTML, stripping tags to get raw text."""
    from html.parser import HTMLParser

    class _TagStripper(HTMLParser):
        def __init__(self):
            super().__init__()
            self._parts: list[str] = []

        def handle_data(self, data: str) -> None:
            self._parts.append(data)

        def get_text(self) -> str:
            return "".join(self._parts)

    raw = path.read_text(encoding="utf-8", errors="replace")
    stripper = _TagStripper()
    stripper.feed(raw)
    text = stripper.get_text()

    return [{
        "text": text,
        "metadata": {
            "source": path.name,
            "file_path": str(path),
            "format": "html",
        },
    }]


def _parse_text(path: Path) -> list[dict]:
    """Parse plain text file."""
    text = path.read_text(encoding="utf-8", errors="replace")
    return [{
        "text": text,
        "metadata": {
            "source": path.name,
            "file_path": str(path),
            "format": "text",
        },
    }]


def _parse_docx(path: Path) -> list[dict]:
    """Extract text from Word .docx including paragraphs and tables."""
    if not _HAS_DOCX:
        logger.warning("python-docx not installed, skipping %s", path.name)
        return []

    doc = DocxDocument(str(path))
    parts = []

    # Paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    # Tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip(" |"):
                parts.append(row_text)

    text = "\n".join(parts)
    return [{
        "text": text,
        "metadata": {
            "source": path.name,
            "file_path": str(path),
            "format": "docx",
        },
    }]


def _parse_xlsx(path: Path) -> list[dict]:
    """Extract text from all sheets in an Excel file."""
    if not _HAS_OPENPYXL:
        logger.warning("openpyxl not installed, skipping %s", path.name)
        return []

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    results = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = " | ".join(cells)
            if line.strip(" |"):
                lines.append(line)
        if lines:
            results.append({
                "text": "\n".join(lines),
                "metadata": {
                    "source": path.name,
                    "sheet": sheet_name,
                    "file_path": str(path),
                    "format": "xlsx",
                },
            })

    wb.close()
    return results


def _parse_csv(path: Path) -> list[dict]:
    """Read CSV file as text, row by row."""
    text_content = path.read_text(encoding="utf-8", errors="replace")
    reader = csv.reader(StringIO(text_content))
    lines = []
    for row in reader:
        line = " | ".join(row)
        if line.strip(" |"):
            lines.append(line)

    return [{
        "text": "\n".join(lines),
        "metadata": {
            "source": path.name,
            "file_path": str(path),
            "format": "csv",
        },
    }]


def _parse_pptx(path: Path) -> list[dict]:
    """Extract text from all slides in a PowerPoint file."""
    if not _HAS_PPTX:
        logger.warning("python-pptx not installed, skipping %s", path.name)
        return []

    prs = Presentation(str(path))
    results = []

    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        parts.append(para.text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        parts.append(row_text)
        if parts:
            results.append({
                "text": "\n".join(parts),
                "metadata": {
                    "source": path.name,
                    "slide": str(i),
                    "file_path": str(path),
                    "format": "pptx",
                },
            })

    return results


def _parse_rtf(path: Path) -> list[dict]:
    """Strip RTF formatting and return plain text."""
    if not _HAS_STRIPRTF:
        logger.warning("striprtf not installed, skipping %s", path.name)
        return []

    raw = path.read_text(encoding="utf-8", errors="replace")
    text = rtf_to_text(raw)

    return [{
        "text": text,
        "metadata": {
            "source": path.name,
            "file_path": str(path),
            "format": "rtf",
        },
    }]


def _parse_json_file(path: Path) -> list[dict]:
    """Flatten a JSON file into text."""
    raw = path.read_text(encoding="utf-8", errors="replace")
    data = json.loads(raw)
    text = _flatten_json(data)

    return [{
        "text": text,
        "metadata": {
            "source": path.name,
            "file_path": str(path),
            "format": "json",
        },
    }]


def _flatten_json(obj, prefix: str = "") -> str:
    """Recursively flatten JSON into key: value lines."""
    lines = []
    if isinstance(obj, dict):
        for key, val in obj.items():
            new_prefix = f"{prefix}.{key}" if prefix else key
            lines.append(_flatten_json(val, new_prefix))
    elif isinstance(obj, list):
        for i, item in enumerate(obj):
            new_prefix = f"{prefix}[{i}]"
            lines.append(_flatten_json(item, new_prefix))
    else:
        return f"{prefix}: {obj}"
    return "\n".join(lines)


def _parse_xml(path: Path) -> list[dict]:
    """Extract text content from XML elements."""
    tree = ET.parse(str(path))
    root = tree.getroot()

    parts = []
    for elem in root.iter():
        if elem.text and elem.text.strip():
            tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
            parts.append(f"{tag}: {elem.text.strip()}")
        if elem.tail and elem.tail.strip():
            parts.append(elem.tail.strip())

    return [{
        "text": "\n".join(parts),
        "metadata": {
            "source": path.name,
            "file_path": str(path),
            "format": "xml",
        },
    }]
